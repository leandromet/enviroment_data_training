from pptx import Presentation
from pptx.util import Inches

# Slide data parsed from the original PDF (condensed and refined from the raw markdown version)
full_slides = [
   
       ("Introdução", [
        "Boas-vindas e objetivos.",
        "Breve visão geral: mudanças climáticas, Saúde Única, SIG e conhecimento tradicional.",
        "Atividade: Apresentação rápida dos participantes."
    ]),
    ("O que é Saúde Única?", [
        "Definição: Integração entre saúde humana, animal e ambiental.",
        "Por que é importante no contexto das mudanças climáticas?"
    ]),
     ("Impactos Interconectados das Mudanças Climáticas", [
        "Saúde Humana: Eventos extremos, mortalidade por calor, doenças infecciosas, deslocamentos.",
        "Meio Ambiente: Florestas afetadas por secas, fontes de água comprometidas.",
        "Ambientes Costeiros: Alterações nas dinâmicas das marés e ondas.",
        "Conclusão: Os impactos são sistêmicos e interligados."
    ]),
    ("Desafios Complexos = Soluções Integradas", [
        "Mudanças climáticas + fatores demográficos, sociais e econômicos.",
        "Ciclos de vetores, patógenos e doenças são multifatoriais.",
        "Saúde Única (One Health): Abordagem abrangente necessária."
    ]),
    ("Vantagens da Abordagem One Health", [
        "Supera limitações das abordagens separadas.",
        "Integra saúde humana, animal e ambiental.",
        "Resultados mais eficazes que métodos convencionais."
    ]),
    ("Estratégias de Adaptação e Mitigação", [
        "Segurança Alimentar: Alimentos de origem animal e pecuária extensiva.",
        "Saneamento Ambiental: Sistemas integrados.",
        "Vigilância Integrada: Resposta coordenada.",
        "Sustentabilidade Social: Políticas éticas e aceitáveis."
    ]),
    ("Vigilância Integrada - Detecção Precoce", [
        "Monitoramento conjunto: humanos + animais.",
        "Identificação precoce de doenças emergentes.",
        "Prevenção antes da transmissão para humanos.",
        "Resultado: Redução significativa de custos."
    ]),
    ("Resultados Holísticos", [
        "Identificação de compensações e sinergias.",
        "Colaboração multissetorial: saúde, agricultura, meio ambiente.",
        "Integração de políticas para resultados eficazes.",
        "Abordagem multidisciplinar."
    ]),
    ("Incorporação de Conhecimentos Tradicionais", [
        "Histórias multigeracionais de povos indígenas.",
        "Conhecimentos locais sobre meio ambiente.",
        "Modelos para compreender variabilidade e mudança.",
        "Integração de saberes tradicionais e científicos."
    ]),

    ("Vantagens da Abordagem Saúde Única", [
        "Comparação com métodos convencionais de saúde pública e animal.",
        "Benefícios para adaptação e mitigação das mudanças climáticas."
    ]),
    ("Desafios de Integração Setorial", [
        "Barreiras para colaboração entre setores.",
        "Importância da comunicação e mecanismos de trabalho conjunto."
    ]),
    ("Atividade Interativa", [
        "Enquete Mentimeter: \"Qual é o maior desafio clima-saúde na sua região?\"",
        "Visual: Nuvem de palavras com respostas."
    ]),
    ("Mudanças Climáticas como Multiplicador de Ameaças", [
        "Conceito de multiplicador de ameaças.",
        "Exemplos gerais de riscos intensificados."
    ]),
    ("Impacto 1 – Zoonoses e Vetores", [
        "Como o desmatamento e mudanças ambientais aumentam o risco de zoonoses.",
        "Exemplos: morcegos, produção pecuária, alternativas à proteína animal."
    ]),
    ("Impacto 2 – Eventos Extremos", [
        "Enchentes, deslocamento de animais, contaminação da água.",
        "Função das áreas úmidas como barreiras naturais."
    ]),
    ("Impacto 2 – Poluição e Bioindicadores", [
        "Macroinvertebrados bentônicos como indicadores de poluição.",
        "Relação entre chuvas intensas e qualidade da água."
    ]),
    ("Impacto 3 – Insegurança Alimentar", [
        "Secas, perda de biodiversidade, serviços ecossistêmicos.",
        "Efeitos do aumento de CO2 em culturas básicas."
    ]),
    ("Impacto 3 – Segurança Alimentar e Pesca", [
        "Queda na pesca e implicações para saúde humana.",
        "Importância de cenários abrangentes para planejamento alimentar."
    ]),
    ("Papel das Comunidades Indígenas", [
        "Detecção precoce de impactos.",
        "Técnicas tradicionais como modelos para políticas públicas."
    ]),
    ("Fatores Espaciais e Sociodemográficos", [
        "Como esses fatores predizem impactos e co-benefícios à saúde.",
        "Exemplos do Brasil e Canadá."
    ]),
    ("Visualização de Hotspots", [
        "Mapa ilustrando hotspots clima-saúde.",
        "Discussão sobre padrões observados."
    ]),
    ("Manejo Cultural do Fogo", [
        "Práticas indígenas na Amazônia.",
        "Prevenção de incêndios e conservação da biodiversidade."
    ]),
    ("SIG Comparativo – Territórios Indígenas", [
        "Imagem SIG: desmatamento em áreas indígenas vs. não indígenas.",
        "Discussão sobre resultados."
    ]),
    ("Discussão em Grupo – Conhecimento Tradicional", [
        "Pergunta: Como incluir conhecimento tradicional no planejamento clima-saúde?",
        "Desafios institucionais e de capacitação."
    ]),
    ("Compartilhamento de Ideias", [
        "Apresentação de propostas dos grupos.",
        "Foto: indígenas no Canadá praticando queimadas controladas."
    ]),
    ("Introdução ao SIG em Saúde Única", [
        "O que é SIG e por que é útil para Saúde Única?",
        "Exemplos de aplicações."
    ]),
    ("SIG – Mapeando Doenças e Clima", [
        "Mapas de expansão de doenças transmitidas por vetores.",
        "Relação entre desmatamento e surtos de malária."
    ]),
    ("SIG – Identificação de Hotspots", [
        "Como identificar áreas de maior risco.",
        "Fatores espaciais em avaliações de impacto à saúde."
    ]),
    ("Demonstração SIG – QGIS", [
        "Projeto exemplo: temperatura vs. distribuição de vetores.",
        "Visual: captura de tela do mapa."
    ]),
    ("Atividade Prática SIG – MapBiomas", [
        "Orientação para uso do MapBiomas.",
        "Exploração de camadas ambientais."
    ]),
    ("Atividade Prática SIG – Global Forest Watch", [
        "Visualização de alertas de desmatamento.",
        "Identificação de padrões e hotspots."
    ]),
    ("Discussão – Aplicações Práticas do SIG", [
        "Como SIG pode informar decisões em saúde animal, pública e ambiental.",
        "Exemplos de uso real."
    ]),
    ("Mudanças de Uso e Cobertura da Terra", [
        "Diferença entre uso e cobertura da terra.",
        "Impactos inesperados do manejo florestal e pastoreio."
    ]),
    ("Casos Reais – Amazônia, Mata Atlântica, Canadá", [
        "Dados MapBiomas: perda de floresta na Amazônia.",
        "Fragmentação da Mata Atlântica e impactos em saúde.",
        "Efeitos do manejo madeireiro no Canadá."
    ]),
    ("Urbanização e Qualidade da Água", [
        "Relação entre urbanização, escoamento e poluição.",
        "Riscos em águas recreativas."
    ]),
    ("Atividade Interativa – Mentimeter", [
        "Qual caso mais te surpreendeu? (Amazônia, Mata Atlântica ou Canadá)",
        "Discussão dos resultados."
    ]),
    ("Cenários Práticos e Encerramento", [
        "Cenário 1: Enchentes e serpentes na Amazônia – plano de intervenção.",
        "Cenário 2: Degelo Ártico e zoonoses – ações prioritárias.",
        "Resumo, recursos e chamada à ação: O que você fará diferente após esta aula?"
    ])
]

speaker_notes = [

    "Dar boas-vindas, apresentar o objetivo da aula e fazer os alunos se apresentarem brevemente.",
    "Explicar o conceito de Saúde Única e sua relevância frente às mudanças climáticas.",
        # 1. Impactos Interconectados das Mudanças Climáticas
    "Os Impactos das Mudanças Climáticas são Interconectados: As mudanças climáticas são consideradas um desafio global urgente, com profundas implicações para a saúde humana. As consequências são amplas e desiguais, incluindo fatalidades e lesões causadas por eventos climáticos extremos, morbidade e mortalidade relacionadas ao calor, aumento da disseminação de doenças infecciosas, deslocamento populacional, perdas econômicas e problemas de saúde mental. Esses impactos vão além da saúde humana, afetando também animais e o meio ambiente. As florestas, por exemplo, provavelmente serão fortemente afetadas pelas mudanças climáticas, o que pode comprometer sua saúde e levar à mortalidade causada por fatores como a seca. Fontes de água, essenciais para a saúde humana, também são impactadas por incêndios, enchentes e infestações de pragas que afetam as florestas, que atuam como purificadores naturais. Ambientes costeiros e a dinâmica microbiana da areia das praias também são influenciados por impactos climáticos, como o aumento da atividade das ondas e das dinâmicas das marés, afetando os riscos à saúde humana.",
    # 2. Desafios Complexos = Soluções Integradas
    "Desafios Integrados Exigem Soluções Integradas: Os efeitos das mudanças climáticas geralmente são difíceis de separar de outras mudanças globais, como alterações demográficas, sociais, econômicas, ambientais e de paisagem. Modificações nos ciclos de vida de vetores, reservatórios e patógenos, bem como doenças de animais domésticos, selvagens e plantas, são influenciadas por múltiplos processos complexos. Enfrentar esses desafios multifacetados requer abordagens integradas. O conceito de Saúde Única (One Health) oferece um contexto mais abrangente para lidar com os efeitos das mudanças climáticas na saúde.",
    # 3. Vantagens da Abordagem One Health
    "Vantagens em Relação a Abordagens Separadas: Uma abordagem integrada de Saúde Única apresenta vantagens claras em relação às abordagens convencionais e separadas de saúde pública e saúde animal para adaptação e mitigação dos efeitos das mudanças climáticas. Essas vantagens vão além da saúde humana e animal.",
    # 4. Estratégias de Adaptação e Mitigação
    "Aprimorando Estratégias de Adaptação e Mitigação: Uma abordagem de Saúde Única para adaptação às mudanças climáticas pode contribuir significativamente em áreas como segurança alimentar, especialmente em relação a alimentos de origem animal e sistemas extensivos de pecuária, saneamento ambiental e desenvolvimento de sistemas integrados de vigilância e resposta. Também auxilia na consideração das implicações sociais mais amplas dos esforços de mitigação climática, buscando garantir que sejam sustentáveis, alinhados com objetivos de saúde pública e aceitáveis ética e publicamente.",
    # 5. Vigilância Integrada - Detecção Precoce
    "Vigilância e Resposta Aprimoradas: Sistemas integrados de vigilância e resposta para humanos e animais são destacados como uma das contribuições mais importantes da abordagem de Saúde Única para mitigar os efeitos das mudanças climáticas. Detectar precocemente doenças emergentes em vetores, animais de criação ou vida selvagem, antes que atinjam humanos, pode evitar custos muito elevados.",
    # 6. Resultados Holísticos
    "Resultados Holísticos: Considerar mudanças climáticas, saúde pública e bem-estar animal sob a ótica da Saúde Única permite identificar compensações e sinergias potenciais, levando a resultados mais holísticos. Ministérios da saúde, agricultura e meio ambiente, juntamente com pesquisadores de diferentes disciplinas, podem colaborar para identificar e avaliar impactos e integrar preocupações em políticas para resultados mais eficazes.",
    # 7. Incorporação de Conhecimentos Tradicionais
    "Incorporação de Conhecimentos Tradicionais: Abordagens integradas podem se beneficiar da incorporação de conhecimentos tradicionais, como as histórias multigeracionais de povos indígenas e comunidades locais em interação com o meio ambiente, que fornecem modelos úteis para compreender variabilidade, incerteza e mudança.",

    "Comparar abordagens tradicionais com Saúde Única, enfatizando benefícios para adaptação.",
    "Apontar obstáculos comuns na colaboração entre setores e a importância da comunicação.",
    "Lançar enquete no Mentimeter sobre desafios locais — mostrar nuvem de palavras.",
    "Apresentar o conceito de 'multiplicador de ameaças' com exemplos atuais.",
    "Mostrar como desmatamento e práticas agrícolas influenciam o surgimento de zoonoses.",
    "Discutir eventos extremos e como afetam ecossistemas e saúde animal/humana.",
    "Apresentar bioindicadores e sua relação com poluição e clima extremo.",
    "Falar sobre impactos climáticos na agricultura e segurança alimentar.",
    "Complementar com dados de pesca e enfatizar importância do planejamento alimentar.",
    "Ressaltar o papel das comunidades indígenas na detecção e resposta a impactos.",
    "Mostrar como fatores geográficos e sociais determinam vulnerabilidades.",
    "Exibir mapa de hotspots e debater padrões visuais e causas subjacentes.",
    "Mostrar práticas indígenas de manejo do fogo como modelo de prevenção.",
    "Comparar áreas indígenas e não indígenas em termos de desmatamento via SIG.",
    "Propor reflexão sobre como integrar o conhecimento tradicional no planejamento.",
    "Apresentar exemplos concretos de propostas feitas por comunidades.",
    "Introduzir o SIG como ferramenta de análise e planejamento em Saúde Única.",
    "Apresentar mapas de vetores e discutir desmatamento e surtos.",
    "Falar sobre métodos de identificação espacial de hotspots com SIG.",
    "Mostrar projeto QGIS com dados climáticos e de vetores — interpretar visual.",
    "Guiar brevemente o uso do MapBiomas — explorar mudanças de cobertura.",
    "Demonstrar alertas do Global Forest Watch e sua interpretação.",
    "Promover debate sobre como SIG pode apoiar decisões intersetoriais.",
    "Explicar diferenças entre uso e cobertura da terra com exemplos.",
    "Apresentar 3 estudos de caso — Amazônia, Mata Atlântica e Canadá.",
    "Debater relação entre urbanização, poluição e saúde em águas recreativas.",
    "Retomar Mentimeter — discutir percepção dos alunos sobre os casos.",
    "Apresentar dois cenários práticos e encerrar com resumo e provocação final."
]

# Create new presentation object
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Curso: Saúde Única, Clima e SIG"
slide.placeholders[1].text = "Mudanças Climáticas, Conhecimento Tradicional e Tecnologias Geoespaciais"

# Generate all slides
for title, bullets in full_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = bullet
            p.level = 0
        else:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0

# Adicionar anotações a cada slide (exceto o título)
for idx, slide in enumerate(list(prs.slides)[1:]):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = speaker_notes[idx]

# Save the full presentation
full_pptx_path = "/media/bndt/db_lin/gitstore/leandromet/enviroment_data_training/onehealth-climate/aula_saude_unica_clima_sig.pptx"
prs.save(full_pptx_path)

full_pptx_path
