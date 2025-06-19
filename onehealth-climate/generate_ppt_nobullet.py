from pptx import Presentation
from pptx.util import Inches

# Slide data parsed from the original PDF (condensed and refined from the raw markdown version)
full_slides = [
    ("Introdução", [
        "Bem-vindos! Nesta aula, vamos abordar a relação entre mudanças climáticas, Saúde Única, SIG e conhecimento tradicional. Começaremos com uma breve apresentação dos participantes."
    ]),
    ("O que é Saúde Única?", [
        "Saúde Única é a integração entre saúde humana, animal e ambiental, fundamental para enfrentar os desafios impostos pelas mudanças climáticas."
    ]),
    ("Impactos Interconectados das Mudanças Climáticas", [
        "As mudanças climáticas afetam a saúde humana, o meio ambiente e ambientes costeiros de forma sistêmica e interligada, exigindo respostas integradas."
    ]),
    ("Desafios Complexos = Soluções Integradas", [
        "Os desafios impostos pelas mudanças climáticas envolvem fatores demográficos, sociais e econômicos, tornando necessária uma abordagem abrangente como a Saúde Única."
    ]),
    ("Vantagens da Abordagem One Health", [
        "A abordagem Saúde Única supera limitações dos métodos convencionais ao integrar saúde humana, animal e ambiental, promovendo resultados mais eficazes."
    ]),
    ("Estratégias de Adaptação e Mitigação", [
        "Adotar Saúde Única contribui para segurança alimentar, saneamento ambiental, vigilância integrada e políticas sustentáveis diante das mudanças climáticas."
    ]),
    ("Vigilância Integrada - Detecção Precoce", [
        "A vigilância integrada permite identificar precocemente doenças emergentes em humanos e animais, reduzindo custos e prevenindo surtos."
    ]),
    ("Resultados Holísticos", [
        "A colaboração entre setores e disciplinas permite identificar sinergias e compensações, resultando em políticas mais eficazes e abrangentes."
    ]),
    ("Incorporação de Conhecimentos Tradicionais", [
        "O conhecimento tradicional de povos indígenas e comunidades locais enriquece a compreensão das mudanças ambientais e contribui para soluções inovadoras."
    ]),
    ("Vantagens da Abordagem Saúde Única", [
        "Comparada a métodos convencionais, Saúde Única oferece benefícios claros para adaptação e mitigação dos efeitos das mudanças climáticas."
    ]),
    ("Desafios de Integração Setorial", [
        "A colaboração entre setores enfrenta barreiras institucionais e de comunicação, mas é essencial para o sucesso das ações integradas."
    ]),
    ("Atividade Interativa", [
        "Vamos realizar uma enquete no Mentimeter para identificar os principais desafios clima-saúde em sua região e visualizar os resultados em uma nuvem de palavras."
    ]),
    ("Mudanças Climáticas como Multiplicador de Ameaças", [
        "As mudanças climáticas intensificam riscos existentes, atuando como multiplicador de ameaças para a saúde e o meio ambiente."
    ]),
    ("Impacto 1 – Zoonoses e Vetores", [
        "O desmatamento e as mudanças ambientais aumentam o risco de zoonoses, como observado em casos envolvendo morcegos e produção pecuária."
    ]),
    ("Impacto 2 – Eventos Extremos", [
        "Eventos extremos, como enchentes, afetam animais e humanos, destacando a importância das áreas úmidas como barreiras naturais."
    ]),
    ("Impacto 2 – Poluição e Bioindicadores", [
        "Macroinvertebrados bentônicos são indicadores de poluição, e chuvas intensas impactam diretamente a qualidade da água."
    ]),
    ("Impacto 3 – Insegurança Alimentar", [
        "Secas e perda de biodiversidade afetam a produção de alimentos, enquanto o aumento de CO2 altera culturas básicas."
    ]),
    ("Impacto 3 – Segurança Alimentar e Pesca", [
        "A redução da pesca impacta a saúde humana e reforça a necessidade de planejamento alimentar abrangente."
    ]),
    ("Papel das Comunidades Indígenas", [
        "Comunidades indígenas desempenham papel fundamental na detecção precoce de impactos e oferecem técnicas tradicionais valiosas para políticas públicas."
    ]),
    ("Fatores Espaciais e Sociodemográficos", [
        "Fatores espaciais e sociodemográficos ajudam a prever impactos e identificar co-benefícios à saúde em diferentes regiões."
    ]),
    ("Visualização de Hotspots", [
        "Mapas de hotspots clima-saúde facilitam a identificação de padrões e a discussão sobre causas subjacentes."
    ]),
    ("Manejo Cultural do Fogo", [
        "Práticas indígenas de manejo do fogo na Amazônia contribuem para a prevenção de incêndios e conservação da biodiversidade."
    ]),
    ("SIG Comparativo – Territórios Indígenas", [
        "Análises SIG mostram diferenças no desmatamento entre áreas indígenas e não indígenas, promovendo discussões sobre conservação."
    ]),
    ("Discussão em Grupo – Conhecimento Tradicional", [
        "Como integrar o conhecimento tradicional no planejamento clima-saúde? Vamos debater desafios institucionais e de capacitação."
    ]),
    ("Compartilhamento de Ideias", [
        "Grupos apresentam propostas e exemplos de integração de saberes, ilustrados por práticas indígenas no Canadá."
    ]),
    ("Introdução ao SIG em Saúde Única", [
        "O SIG é uma ferramenta essencial para análise e planejamento em Saúde Única, permitindo aplicações inovadoras."
    ]),
    ("SIG – Mapeando Doenças e Clima", [
        "Mapas SIG revelam a expansão de doenças transmitidas por vetores e a relação com o desmatamento."
    ]),
    ("SIG – Identificação de Hotspots", [
        "O SIG permite identificar áreas de maior risco e avaliar fatores espaciais em impactos à saúde."
    ]),
    ("Demonstração SIG – QGIS", [
        "Exemplo prático: análise da relação entre temperatura e distribuição de vetores usando QGIS."
    ]),
    ("Atividade Prática SIG – MapBiomas", [
        "Vamos explorar o MapBiomas para analisar camadas ambientais e mudanças de cobertura da terra."
    ]),
    ("Atividade Prática SIG – Global Forest Watch", [
        "Visualizaremos alertas de desmatamento no Global Forest Watch e identificaremos padrões relevantes."
    ]),
    ("Discussão – Aplicações Práticas do SIG", [
        "O SIG apoia decisões em saúde animal, pública e ambiental, com exemplos reais de aplicação."
    ]),
    ("Mudanças de Uso e Cobertura da Terra", [
        "Entenda a diferença entre uso e cobertura da terra e os impactos inesperados do manejo florestal e pastoreio."
    ]),
    ("Casos Reais – Amazônia, Mata Atlântica, Canadá", [
        "Estudos de caso mostram a perda de floresta na Amazônia, fragmentação da Mata Atlântica e efeitos do manejo madeireiro no Canadá."
    ]),
    ("Urbanização e Qualidade da Água", [
        "A urbanização influencia o escoamento e a poluição, aumentando riscos em águas recreativas."
    ]),
    ("Atividade Interativa – Mentimeter", [
        "Qual caso mais te surpreendeu? Vamos discutir os resultados da enquete sobre Amazônia, Mata Atlântica e Canadá."
    ]),
    ("Cenários Práticos e Encerramento", [
        "Vamos analisar dois cenários práticos e encerrar com um resumo e uma provocação final sobre ações futuras."
    ])
]

speaker_notes = [

    "Dar boas-vindas, apresentar o objetivo da aula e fazer os alunos se apresentarem brevemente.",
    "Explicar o conceito de Saúde Única e sua relevância frente às mudanças climáticas.",
        # 1. Impactos Interconectados das Mudanças Climáticas
    "Os Impactos das Mudanças Climáticas são Interconectados: As mudanças climáticas são consideradas um desafio global urgente, com profundas implicações para a saúde humana. As consequências são amplas e desiguais, incluindo fatalidades e lesões causadas por eventos climáticos extremos, morbidade e mortalidade relacionadas ao calor, aumento da disseminação de doenças infecciosas, deslocamento populacional, perdas econômicas e problemas de saúde mental. Esses impactos vão além da saúde humana, afetando também animais e o meio ambiente. As florestas, por exemplo, provavelmente serão fortemente afetadas pelas mudanças climáticas, o que pode comprometer sua saúde e levar à mortalidade causada por fatores como a seca. Fontes de água, essenciais para a saúde humana, também são impactadas por incêndios, enchentes e infestações de pragas que afetam as florestas, que atuam como purificadores naturais. Ambientes costeiros e a dinâmica microbiana da areia das praias também são influenciados por impactos climáticos, como o aumento da atividade das ondas e das dinâmicas das marés, afetando os riscos à saúde humana.",
    # 2. Desafios Complexos = Soluções Integradas
    "Desafios Integrados Exigem Soluções Integradas: Os efeitos das mudanças climáticas geralmente são difíceis de separar de outras mudanças globais, como alterações demográficas, sociais, econômicas, ambientais e de paisagem. Modificações nos ciclos de vida de vetores, reservatórios e patógenos, bem como doenças de animais domésticos, selvagens e plantas, são influenciadas por múltiplos processos complexos. Enfrentar esses desafios multifacetados requer abordagens integradas. O conceito de Saúde Única (One Health) oferece um contexto mais abrangente para lidar com os efeitos das mudanças climáticas na saúde.",
    # 3. Vantagens da Abordagem One Health
    "Vantagens em Relação a Abordagens Separadas: Uma abordagem integrada de Saúde Única apresenta vantagens claras em relação às abordagens convencionais e separadas de saúde pública e saúde animal para adaptação e mitigação dos efeitos das mudanças climáticas. Essas vantagens vão além da saúde humana e animal.",
    # 4. Estratégias de Adaptação e Mitigação
    "Aprimorando Estratégias de Adaptação e Mitigação: Uma abordagem de Saúde Única para adaptação às mudanças climáticas pode contribuir significativamente em áreas como segurança alimentar, especialmente em relação a alimentos de origem animal e sistemas extensivos de pecuária, saneamento ambiental e desenvolvimento de sistemas integrados de vigilância e resposta. Também auxilia na consideração das implicações sociais mais amplas dos esforços de mitigação climática, buscando garantir que sejam sustentáveis, alinhados com objetivos de saúde pública e aceitáveis ética e publicamente.",
    # 5. Vigilância Integrada - Detecção Precoce
    "Vigilância e Resposta Aprimoradas: Sistemas integrados de vigilância e resposta para humanos e animais são destacados como uma das contribuições mais importantes da abordagem de Saúde Única para mitigar os efeitos das mudanças climáticas. Detectar precocemente doenças emergentes em vetores, animais de criação ou vida selvagem, antes que atinjam humanos, pode evitar custos muito elevados.",
    # 6. Resultados Holísticos
    "Resultados Holísticos: Considerar mudanças climáticas, saúde pública e bem-estar animal sob a ótica da Saúde Única permite identificar compensações e sinergias potenciais, levando a resultados mais holísticos. Ministérios da saúde, agricultura e meio ambiente, juntamente com pesquisadores de diferentes disciplinas, podem colaborar para identificar e avaliar impactos e integrar preocupações em políticas para resultados mais eficazes.",
    # 7. Incorporação de Conhecimentos Tradicionais
    "Incorporação de Conhecimentos Tradicionais: Abordagens integradas podem se beneficiar da incorporação de conhecimentos tradicionais, como as histórias multigeracionais de povos indígenas e comunidades locais em interação com o meio ambiente, que fornecem modelos úteis para compreender variabilidade, incerteza e mudança.",

    "Comparar abordagens tradicionais com Saúde Única, enfatizando benefícios para adaptação.",
    "Apontar obstáculos comuns na colaboração entre setores e a importância da comunicação.",
    "Lançar enquete no Mentimeter sobre desafios locais — mostrar nuvem de palavras.",
    "Apresentar o conceito de 'multiplicador de ameaças' com exemplos atuais.",
    "Mostrar como desmatamento e práticas agrícolas influenciam o surgimento de zoonoses.",
    "Discutir eventos extremos e como afetam ecossistemas e saúde animal/humana.",
    "Apresentar bioindicadores e sua relação com poluição e clima extremo.",
    "Falar sobre impactos climáticos na agricultura e segurança alimentar.",
    "Complementar com dados de pesca e enfatizar importância do planejamento alimentar.",
    "Ressaltar o papel das comunidades indígenas na detecção e resposta a impactos.",
    "Mostrar como fatores geográficos e sociais determinam vulnerabilidades.",
    "Exibir mapa de hotspots e debater padrões visuais e causas subjacentes.",
    "Mostrar práticas indígenas de manejo do fogo como modelo de prevenção.",
    "Comparar áreas indígenas e não indígenas em termos de desmatamento via SIG.",
    "Propor reflexão sobre como integrar o conhecimento tradicional no planejamento.",
    "Apresentar exemplos concretos de propostas feitas por comunidades.",
    "Introduzir o SIG como ferramenta de análise e planejamento em Saúde Única.",
    "Apresentar mapas de vetores e discutir desmatamento e surtos.",
    "Falar sobre métodos de identificação espacial de hotspots com SIG.",
    "Mostrar projeto QGIS com dados climáticos e de vetores — interpretar visual.",
    "Guiar brevemente o uso do MapBiomas — explorar mudanças de cobertura.",
    "Demonstrar alertas do Global Forest Watch e sua interpretação.",
    "Promover debate sobre como SIG pode apoiar decisões intersetoriais.",
    "Explicar diferenças entre uso e cobertura da terra com exemplos.",
    "Apresentar 3 estudos de caso — Amazônia, Mata Atlântica e Canadá.",
    "Debater relação entre urbanização, poluição e saúde em águas recreativas.",
    "Retomar Mentimeter — discutir percepção dos alunos sobre os casos.",
    "Apresentar dois cenários práticos e encerrar com resumo e provocação final."
]

# Create new presentation object
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Curso: Saúde Única, Clima e SIG"
slide.placeholders[1].text = "Mudanças Climáticas, Conhecimento Tradicional e Tecnologias Geoespaciais"

# Generate all slides
for title, bullets in full_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = bullet
            p.level = 0
        else:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0

# Adicionar anotações a cada slide (exceto o título)
for idx, slide in enumerate(list(prs.slides)[1:]):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = speaker_notes[idx]

# Save the full presentation
full_pptx_path = "/media/bndt/db_lin/gitstore/leandromet/enviroment_data_training/onehealth-climate/aula_saude_unica_clima_sig_v3.pptx"
prs.save(full_pptx_path)

full_pptx_path
